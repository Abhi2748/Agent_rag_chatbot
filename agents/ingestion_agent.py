import os
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

class IngestionAgent:
    SUPPORTED_FORMATS = ['.pdf', '.docx', '.pptx', '.csv', '.txt', '.md']

    @staticmethod
    def process(files, trace_id):
        chunks = []
        for file in files:
            filename = getattr(file, 'name', None) or getattr(file, 'filename', 'unknown')
            ext = os.path.splitext(filename)[-1].lower()

            if ext not in IngestionAgent.SUPPORTED_FORMATS:
                continue

            file.seek(0)  # Reset pointer before reading
            text = ''
            try:
                if ext == '.pdf':
                    reader = PdfReader(file)
                    text = "\n".join([page.extract_text() or '' for page in reader.pages])

                elif ext == '.docx':
                    doc = Document(file)
                    text = "\n".join([p.text for p in doc.paragraphs])

                elif ext == '.pptx':
                    prs = Presentation(file)
                    text = "\n".join([shape.text for slide in prs.slides
                                      for shape in slide.shapes if hasattr(shape, 'text')])

                elif ext == '.csv':
                    df = pd.read_csv(file)
                    text = df.to_csv(index=False)

                elif ext in ['.txt', '.md']:
                    try:
                        text = file.read().decode('utf-8')
                    except UnicodeDecodeError:
                        file.seek(0)
                        text = file.read().decode('latin1')

            except Exception as e:
                print(f"[ERROR] Failed to extract from {filename}: {e}")
                continue

            print(f"Extracted from {filename}: {repr(text[:100])}...")

            # Simple chunking
            for i in range(0, len(text), 500):
                chunk = text[i:i+500]
                if chunk.strip():
                    chunks.append({
                        'chunk': chunk,
                        'source': filename,
                        'trace_id': trace_id
                    })

        return chunks
